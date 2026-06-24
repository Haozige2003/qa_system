import os
os.environ["HF_ENDPOINT"] = "https://hf-mirror.com"
import gradio as gr
import os
import asyncio
import time
import re
from openai import OpenAI
from rag_retriever import RAGRetriever
from kreuzberg import extract_file

retriever = RAGRetriever(persist_dir="./chroma_db")

DASHSCOPE_API_KEY = os.environ.get("DASHSCOPE_API_KEY", "sk-xxx")
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "sk-yyy")

API_ENDPOINTS = [
    {"name": "DeepSeek", "base_url": "https://api.deepseek.com/v1", "api_key": DEEPSEEK_API_KEY,
     "model": "deepseek-chat"},
    {"name": "阿里云", "base_url": "https://dashscope.aliyuncs.com/compatible-mode/v1", "api_key": DASHSCOPE_API_KEY,
     "model": "deepseek-r1"}
]

def call_llm(prompt):
    for ep in API_ENDPOINTS:
        if not ep["api_key"]:
            continue
        for attempt in range(2):
            try:
                c = OpenAI(api_key=ep["api_key"], base_url=ep["base_url"], timeout=60)
                resp = c.chat.completions.create(
                    model=ep["model"],
                    messages=[{"role": "system", "content": "根据事实回答"}, {"role": "user", "content": prompt}],
                    temperature=0.3, max_tokens=500
                )
                return resp.choices[0].message.content.strip()
            except:
                time.sleep(2 ** attempt)
    return "API调用失败"

# 4. 读取文件
def read_file(file):
    if not file:
        return ""
    ext = os.path.splitext(file.name)[1].lower()
    if ext in ['.txt', '.md', '.csv', '.json', '.py']:
        for enc in ['utf-8', 'gbk']:
            try:
                with open(file.name, 'r', encoding=enc) as f:
                    return re.sub(r'\s+', ' ', f.read()).strip()
            except:
                continue
    try:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        res = loop.run_until_complete(extract_file(file.name))
        loop.close()
        if res.content and res.content.strip():
            return re.sub(r'\s+', ' ', res.content).strip()
    except:
        pass
    try:
        if ext == '.pdf':
            try:
                import pypdfium2 as pdfium
                pdf = pdfium.PdfDocument(file.name)
                text = "".join([page.get_text() for page in pdf])
                pdf.close()
                if text.strip():
                    return re.sub(r'\s+', ' ', text).strip()
            except:
                pass
            import PyPDF2
            with open(file.name, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text = "".join([p.extract_text() or "" for p in reader.pages])
                if text.strip():
                    return re.sub(r'\s+', ' ', text).strip()
        elif ext == '.docx':
            import docx
            doc = docx.Document(file.name)
            text = "".join([p.text for p in doc.paragraphs if p.text.strip()])
            if text.strip():
                return re.sub(r'\s+', ' ', text).strip()
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file.name)
            text = "".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
            if text.strip():
                return re.sub(r'\s+', ' ', text).strip()
        elif ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(file.name, read_only=True, data_only=True)
            rows = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = " ".join([str(c) for c in row if c])
                    if row_text:
                        rows.append(row_text)
            text = "".join(rows)
            if text.strip():
                return re.sub(r'\s+', ' ', text).strip()
    except:
        pass
    return ""

# 5. 文件入库（优化：保留原始文件名作为 doc_id，不强行改为“X简历”）
def add_files(files):
    if not files:
        total = retriever.count_documents()
        return f"未选择文件，当前知识库总块数：{total}"

    logs, ok = [], 0
    for f in files:
        content = read_file(f)
        if not content or len(content) < 20:
            logs.append(f" {os.path.basename(f.name)} 解析失败")
            continue

        # 提取人名（仅用于元数据过滤，不强制改变文档 ID）
        name = retriever.extract_person_name_from_text(content)

        #doc_id 始终保留原始文件名（不含后缀）
        doc_id = os.path.splitext(os.path.basename(f.name))[0]

        retriever.delete_document(doc_id)

        cnt = retriever.add_document(content, doc_id=doc_id, person_name=name, chunk_size=500, overlap=150)
        ok += 1
        logs.append(f" {os.path.basename(f.name)} → {cnt}块" + (f" (关联人名：{name})" if name else ""))

    total = retriever.count_documents()
    # 将总块数整合进日志，只返回一个字符串，避免界面被数字覆盖
    return f"入库完成：成功 {ok} 个文件，总块数 {total}\n" + "\n".join(logs)

# 6. 问答
def ask(question):
    if not question:
        return "请输入问题", ""

    name = retriever.extract_person_name_from_question(question)
    docs = []
    if name:
        # 按人名过滤检索，两个文档都能被搜到
        docs = retriever.retrieve(question, person_name=name, top_k=6)  # 适当增加检索数量
        if len(docs) < 2:
            extra = retriever.retrieve(question, top_k=6)
            seen = set([d[0] for d in docs])
            for d in extra:
                if d[0] not in seen:
                    docs.append(d)
                    seen.add(d[0])
    else:
        docs = retriever.retrieve(question, top_k=6)

    if not docs:
        return "未找到相关信息", ""

    # 构建上下文（明确标注文档ID）
    context = "\n\n".join([f"【文档来源：{doc_id}】\n{text}" for doc_id, text in docs])

    prompt = f"""
你是一个严谨的文档对比分析助手。下面提供了来自不同文档的片段，每个片段前面都有【文档来源：XXX】的标记。

请根据用户的问题，执行以下操作：
1. 如果问题涉及比较（如“区别”、“差异”、“哪个更好”、"几者间的联系"），请分别引用不同【文档来源】的内容，明确指出它们之间的相同点和不同点。
2. 如果问题不涉及比较，则正常根据提供的信息回答。
3. 严禁使用你自己记忆中的知识，只基于下面提供的片段回答。

提供的片段：
{context}

用户问题：{question}
"""
    ans = call_llm(prompt)
    return ans, context

def refresh():
    total = retriever.count_documents()
    ids = retriever.get_all_doc_ids()
    return f"知识库状态：{total}块，文档ID列表：{', '.join(ids) if ids else '空'}"

def clear():
    retriever.clear_all()
    return " 知识库已完全清空"

with gr.Blocks() as demo:
    gr.Markdown("#  文档问答系统（多格式支持）")
    with gr.Row():
        with gr.Column():
            files = gr.File(label="📎 上传文件（可多选）", file_count="multiple")
            btn_up = gr.Button(" 入库")
            btn_clr = gr.Button(" 清空")
            btn_ref = gr.Button(" 刷新")
            status = gr.Textbox(label=" 状态日志", lines=8)
        with gr.Column():
            question = gr.Textbox(label=" 输入您的问题")
            btn_ask = gr.Button(" 回答")
            answer = gr.Textbox(label=" 答案", lines=6)
            context = gr.Textbox(label=" 检索到的原文片段", lines=8)

    # 绑定事件：所有回调函数现在都只返回一个字符串，统一更新 status
    btn_up.click(add_files, [files], status)
    btn_clr.click(clear, None, status)
    btn_ref.click(refresh, None, status)
    btn_ask.click(ask, [question], [answer, context])

if __name__ == "__main__":
    demo.launch()
